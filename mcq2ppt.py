from pptx import Presentation
from pptx.util import Inches

def create_ppt(mcq_data, output_file):
    # Create a presentation object
    prs = Presentation()

    for idx, question_data in enumerate(mcq_data):
        question = question_data['question']
        options = question_data['options']
        answer = question_data['answer']

        # Create a slide for each question
        slide_layout = prs.slide_layouts[1]  # Use the 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)

        # Set the title to the question
        title = slide.shapes.title
        title.text = f"Q{idx + 1}: {question}"

        # Add the options and answer as bullet points
        content = slide.placeholders[1].text_frame
        for option in options:
            p = content.add_paragraph()
            p.text = option

        # Create a slide for the answer
        slide_layout = prs.slide_layouts[1]  # Use the 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)
        # Set the title to "Solution"
        title = slide.shapes.title
        title.text = "Solution"
        # Add the answer as a bullet point
        content = slide.placeholders[1].text_frame
        p = content.add_paragraph()
        p.text = answer

    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example MCQ data
mcq_data = [
    {"question": "Which of the following is used to terminate a loop in C?",
     "options": ["A. break", "B. continue", "C. exit", "D. stop"],
     "answer": "A. break"},
    {"question": "What is the purpose of the 'default' case in a switch statement?",
     "options": ["A. To specify the starting point", "B. To handle unspecified cases", "C. To end the switch statement", "D. To restart the switch statement"],
     "answer": "B. To handle unspecified cases"}
]

# Create the PowerPoint presentation
create_ppt(mcq_data, "mcq_presentation.pptx")
